"""
Step 6: Final Deliverables
============================
- Comprehensive Excel workbook
- PowerPoint presentation
- Summary Word report
"""
import os
import pandas as pd
import numpy as np
from collections import Counter
from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from docx import Document
from docx.shared import Inches as DocInches, Pt as DocPt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from config import (
    ENCOUNTER_CLEAN_PATH, ENCOUNTER_PRESCRIBER_PATH,
    ENCOUNTER_CLUSTERED_PATH, FIGURES_DIR, REPORTS_DIR,
)


def create_excel_workbook(df_clean, df_prescriber, df_clustered):
    """Create comprehensive Excel workbook with multiple sheets."""
    print("\n--- Creating Excel workbook ---")
    path = f"{REPORTS_DIR}/final_deliverables.xlsx"

    with pd.ExcelWriter(path, engine="openpyxl") as writer:
        # Sheet 1: Cluster summary
        clusters = sorted(df_clustered["cluster"].unique())
        summary_rows = []
        for c in clusters:
            sub = df_clustered[df_clustered["cluster"] == c]
            summary_rows.append({
                "Cluster": c,
                "N": len(sub),
                "Pct": round(len(sub) / len(df_clustered) * 100, 1),
                "AB_Rate_%": round(sub["has_antibiotic"].mean() * 100, 1),
                "Mono_%": round(sub["antibiotic_monotherapy"].mean() * 100, 1),
                "Combo_%": round(sub["antibiotic_combination"].mean() * 100, 1),
                "Polypharm_%": round(sub["polypharmacy_flag"].mean() * 100, 1),
                "Mean_Drugs": round(sub["num_distinct_drugs"].mean(), 2),
                "Mean_Age_Months": round(sub["age_months"].mean(), 1),
                "Dominant_Age": sub["age_stratum"].mode().iloc[0] if len(sub) > 0 else "",
            })
        pd.DataFrame(summary_rows).to_excel(
            writer, sheet_name="Cluster Summary", index=False
        )

        # Sheet 2: Full data with clusters (limit rows for Excel)
        max_rows = min(len(df_clustered), 100000)
        df_clustered.head(max_rows).to_excel(
            writer, sheet_name="Encounter Data", index=False
        )

        # Sheet 3: Prescriber analysis
        prescriber_stats = df_clustered.groupby("prescriber_id").agg(
            total_encounters=("OPDID", "count"),
            ab_rate=("has_antibiotic", "mean"),
            polypharmacy_rate=("polypharmacy_flag", "mean"),
            mean_drugs=("num_distinct_drugs", "mean"),
            clusters_seen=("cluster", lambda x: ", ".join(map(str, sorted(x.unique())))),
        ).reset_index()
        prescriber_stats["ab_rate"] = (prescriber_stats["ab_rate"] * 100).round(1)
        prescriber_stats["polypharmacy_rate"] = (prescriber_stats["polypharmacy_rate"] * 100).round(1)
        prescriber_stats.sort_values("total_encounters", ascending=False).to_excel(
            writer, sheet_name="Prescriber Analysis", index=False
        )

        # Sheet 4: Top drugs per cluster
        drug_rows = []
        for c in clusters:
            sub = df_clustered[df_clustered["cluster"] == c]
            drugs = []
            for drug_str in sub["drug_names"].dropna():
                for d in str(drug_str).split("|"):
                    d = d.strip()
                    if d and d.lower() != "nan":
                        drugs.append(d)
            for drug, count in Counter(drugs).most_common(15):
                drug_rows.append({"Cluster": c, "Drug": drug, "Count": count})
        pd.DataFrame(drug_rows).to_excel(
            writer, sheet_name="Top Drugs by Cluster", index=False
        )

        # Sheet 5: Validation results (if exists)
        val_path = f"{REPORTS_DIR}/validation_results.xlsx"
        if os.path.exists(val_path):
            val_df = pd.read_excel(val_path, sheet_name="Bootstrap Summary")
            val_df.to_excel(writer, sheet_name="Validation", index=False)

    print(f"  Excel workbook saved: {path}")


def create_powerpoint(df_clustered):
    """Create PowerPoint presentation."""
    print("\n--- Creating PowerPoint presentation ---")
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Antibiotic Prescribing Archetypes"
    slide.placeholders[1].text = (
        "Clustering Analysis of Paediatric OPD Prescriptions\n"
        "January 2026 - Lady Ridgeway Hospital\n"
        "MD Health Informatics Research"
    )

    # Slide 2: Data overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Data Overview"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = f"Total encounters: {len(df_clustered):,}"
    tf.add_paragraph().text = f"Unique patients: {df_clustered['patient_id'].nunique():,}"
    tf.add_paragraph().text = f"Unique prescribers: {df_clustered['prescriber_id'].nunique():,}"
    tf.add_paragraph().text = f"Date range: January 2026"
    tf.add_paragraph().text = f"Antibiotic rate: {df_clustered['has_antibiotic'].mean()*100:.1f}%"
    tf.add_paragraph().text = f"Polypharmacy rate: {df_clustered['polypharmacy_flag'].mean()*100:.1f}%"

    # Slide 3: Methodology
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Methodology"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "K-medoids clustering with Gower distance"
    tf.add_paragraph().text = "Features: age, weight, drugs, antibiotics, prescriber rates"
    tf.add_paragraph().text = "Mixed data types: numerical (scaled), binary, categorical (one-hot)"
    tf.add_paragraph().text = f"Optimal k selected via silhouette score (k={df_clustered['cluster'].nunique()})"
    tf.add_paragraph().text = "Validated with bootstrap resampling (100 iterations)"

    # Slides for each cluster
    clusters = sorted(df_clustered["cluster"].unique())
    for c in clusters:
        sub = df_clustered[df_clustered["cluster"] == c]
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = f"Cluster {c}"
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = f"Size: {len(sub):,} encounters ({len(sub)/len(df_clustered)*100:.1f}%)"
        tf.add_paragraph().text = f"Dominant age: {sub['age_stratum'].mode().iloc[0]}"
        tf.add_paragraph().text = f"Antibiotic rate: {sub['has_antibiotic'].mean()*100:.1f}%"
        tf.add_paragraph().text = f"Polypharmacy: {sub['polypharmacy_flag'].mean()*100:.1f}%"
        tf.add_paragraph().text = f"Mean drugs: {sub['num_distinct_drugs'].mean():.2f}"
        tf.add_paragraph().text = f"Combination AB: {sub['antibiotic_combination'].mean()*100:.1f}%"

    # Add figure slides
    figure_files = [
        ("10_pca_clusters.png", "PCA Cluster Visualization"),
        ("11_cluster_heatmap.png", "Cluster Characteristics Heatmap"),
        ("12_cluster_comparison_bars.png", "Cluster Comparison"),
        ("09_silhouette_scores.png", "Silhouette Score Analysis"),
    ]
    for fig_file, title in figure_files:
        fig_path = os.path.join(FIGURES_DIR, fig_file)
        if os.path.exists(fig_path):
            slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
            slide.shapes.add_picture(fig_path, Inches(1), Inches(0.5),
                                     width=Inches(11), height=Inches(6.5))

    path = f"{REPORTS_DIR}/final_presentation.pptx"
    prs.save(path)
    print(f"  PowerPoint saved: {path}")


def create_summary_report(df_clustered):
    """Create summary Word report."""
    print("\n--- Creating summary Word report ---")
    doc = Document()

    doc.add_heading("Summary Report: Antibiotic Prescribing Clustering Analysis", level=0)
    doc.add_paragraph(
        "Paediatric OPD, Lady Ridgeway Hospital\n"
        "January 2026 | MD Health Informatics Research"
    )

    # Methods
    doc.add_heading("Methods", level=1)
    doc.add_paragraph(
        "Encounter-level data was extracted from HHIMS for January 2026. "
        "Each encounter was characterised by patient demographics, prescribing "
        "patterns, and prescriber behaviour metrics. K-medoids clustering with "
        "Gower distance was applied to identify distinct prescribing archetypes. "
        "Cluster validity was assessed using silhouette scores and bootstrap "
        "resampling (100 iterations, 80% samples)."
    )

    # Descriptive statistics
    doc.add_heading("Descriptive Statistics", level=1)
    doc.add_paragraph(f"Total encounters: {len(df_clustered):,}")
    doc.add_paragraph(f"Unique patients: {df_clustered['patient_id'].nunique():,}")
    doc.add_paragraph(f"Unique prescribers: {df_clustered['prescriber_id'].nunique():,}")
    doc.add_paragraph(f"Overall antibiotic rate: {df_clustered['has_antibiotic'].mean()*100:.1f}%")
    doc.add_paragraph(f"Polypharmacy rate: {df_clustered['polypharmacy_flag'].mean()*100:.1f}%")

    # Cluster results
    doc.add_heading("Cluster Results", level=1)
    n_clusters = df_clustered["cluster"].nunique()
    doc.add_paragraph(f"Number of clusters identified: {n_clusters}")

    for c in sorted(df_clustered["cluster"].unique()):
        sub = df_clustered[df_clustered["cluster"] == c]
        doc.add_heading(f"Cluster {c}", level=2)
        doc.add_paragraph(f"Size: {len(sub):,} ({len(sub)/len(df_clustered)*100:.1f}%)")
        doc.add_paragraph(f"Antibiotic rate: {sub['has_antibiotic'].mean()*100:.1f}%")
        doc.add_paragraph(f"Polypharmacy rate: {sub['polypharmacy_flag'].mean()*100:.1f}%")
        doc.add_paragraph(f"Mean drugs per encounter: {sub['num_distinct_drugs'].mean():.2f}")
        doc.add_paragraph(f"Dominant age group: {sub['age_stratum'].mode().iloc[0]}")

    # Add figures
    doc.add_heading("Figures", level=1)
    for fig in ["10_pca_clusters.png", "11_cluster_heatmap.png",
                "12_cluster_comparison_bars.png", "13_cluster_sizes.png"]:
        fig_path = os.path.join(FIGURES_DIR, fig)
        if os.path.exists(fig_path):
            doc.add_picture(fig_path, width=DocInches(5.5))
            doc.add_paragraph("")

    path = f"{REPORTS_DIR}/summary_report.docx"
    doc.save(path)
    print(f"  Summary report saved: {path}")


def run():
    print("=" * 60)
    print("STEP 6: FINAL DELIVERABLES")
    print("=" * 60)

    df_clean = pd.read_csv(ENCOUNTER_CLEAN_PATH, low_memory=False)
    df_prescriber = pd.read_csv(ENCOUNTER_PRESCRIBER_PATH, low_memory=False)
    df_clustered = pd.read_csv(ENCOUNTER_CLUSTERED_PATH, low_memory=False)
    print(f"\nLoaded all datasets. Clustered: {len(df_clustered)} encounters")

    create_excel_workbook(df_clean, df_prescriber, df_clustered)
    create_powerpoint(df_clustered)
    create_summary_report(df_clustered)

    print("\n" + "=" * 60)
    print("ALL DELIVERABLES COMPLETE!")
    print("=" * 60)


if __name__ == "__main__":
    run()
